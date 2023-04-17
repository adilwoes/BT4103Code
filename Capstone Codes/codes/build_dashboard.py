import re
import warnings
warnings.filterwarnings("ignore")

from ipywidgets import Button, Layout, HBox, \
        VBox, Tab
import plotly.graph_objects as go

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from codes.tool_util_viz import ToolUtilViz
from codes.technode_viz import TechNodeViz
from codes.bu_viz import BUViz
from codes.Anomaly_Detection import Anomaly_Detection
from codes.TAT_QBreakdown_dashboard import TatQBreakdown

files = {
    "Data/M4_Utilization_Sheet.txt": "M4",
    "Data/M5_Utilization_Sheet.txt": "M5",
    "Data/CP_Utilization_Sheet.txt": "CP"
}

#file_names = [f for f in os.listdir(directory) if f.endswith('.xlsm')]
file_names = ['Data/Singapore_Device_Priority_2022 - WW09']
analyse_year = re.findall('\d+', file_names[0])[0]

excel_file = 'Data/Job Input Form_cleaned.xlsm'
anomaly_detection_file_path = 'Data/tech_anomaly_detection_rate.xlsx'
tat_qbreakdown_file_path = f'Data/Singapore_Device_Priority_{analyse_year} - Calculated.xlsx'
tat_qbreakdown_file_path_cancelled = f'Data/Singapore_Device_Priority_{analyse_year} - Cancelled.xlsx'
ww = 'Data/Work Week Calendar.xlsx'

class BuildDashboard:
    """
    A class for building a dashboard.

    Attributes
    ----------
    TAB_TITLES : list
        A list of tab titles for the dashboard. 
    tat_qbreakdown : TatQBreakdown
        An instance of TatQBreakdown.
    anomaly_detection : Anomaly_Detection
        An instance of Anomaly_Detection.
    tool_util_viz : ToolUtilViz
        An instance of ToolUtilViz.
    tech_node_viz : TechNodeViz
        An instance of TechNodeViz.
    bu_viz : BUViz
        An instance of BUViz.
    export_to_ppt : Button
        A button to export KPI graphs to PPT.
    export_to_ppt1 : Button
        A button to export all graphs to PPT.
    dashboard : Tab
        A Tab instance for the dashboard.
    export : HBox
        An HBox instance for the export.
    app : VBox
        A VBox instance for the app.
    """

    TAB_TITLES = ['KPI', 'KPI (others)', 'BU Loading', 'Tool Utilization', \
                  'Technology Nodes', 'Technology Anomaly Detection Rate']

    def __init__(self):
        self.tat_qbreakdown = TatQBreakdown(tat_qbreakdown_file_path, tat_qbreakdown_file_path_cancelled, ww)
        self.anomaly_detection = Anomaly_Detection(anomaly_detection_file_path)
        self.tool_util_viz = ToolUtilViz(files, ww)
        self.tech_node_viz = TechNodeViz(excel_file)
        self.bu_viz = BUViz(excel_file)
        self.export_to_ppt = Button(description='Export KPI Graphs to PPT', button_style='warning', \
                                    layout = Layout(width='auto'))
        self.export_to_ppt1 = Button(description='Export All Graphs to PPT', button_style='warning', \
                                     layout = Layout(width='auto'))
        self.export_to_ppt.on_click(self.create_ppt)
        self.export_to_ppt1.on_click(self.create_ppt_all)
        
        self.dashboard = Tab([self.tat_qbreakdown.kpi_tab, self.tat_qbreakdown.kpi_others_tab, \
               self.bu_viz.tab, self.tool_util_viz.tab, self.tech_node_viz.tab, self.anomaly_detection.tab])
        [self.dashboard.set_title(i, title) for i, title in enumerate(self.TAB_TITLES)]
        
        self.export = HBox([self.export_to_ppt, self.export_to_ppt1])
        self.app = VBox([self.export, self.dashboard])

    def save_png_tab1(self):
        """
        Saves the images of various graphs and visualizations displayed in the tabs of the GUI.
        This method retrieves the various graph and visualization objects from the GUI tabs and saves them as PNG images in a
        new 'Output' directory. The saved images include the Tool Utilization, BU Loading, TAT, High Priority Queue Time
        Breakdown, Overall Analysis Status, Type Breakdown, Type Breakdown by Month, Product Breakdown, Technology Node
        Distribution, Technology Product Distribution, and Technology Failure Distribution graphs and visualizations.

        Parameters:
		-----------
		None
		
		Returns:
		--------
		None
        """
        turnaround_graph = self.tat_qbreakdown.kpi_tab.children[5]
        priority_graph = self.tat_qbreakdown.kpi_tab.children[8]
        csc_graph = self.tat_qbreakdown.kpi_tab.children[10]
        bu_graph = self.bu_viz.tab.children[2]
        tool_util = self.tool_util_viz.tab.children[2]
        
        type_graph = self.tat_qbreakdown.kpi_others_tab.children[5]
        type_month_graph = self.tat_qbreakdown.kpi_others_tab.children[7]
        prod_graph = self.tat_qbreakdown.kpi_others_tab.children[9]
        
        tech_node = self.tech_node_viz.tab.children[1].children[0]
        prod_dist = self.tech_node_viz.tab.children[3].children[0].children[1]
        failure_dist = self.tech_node_viz.tab.children[3].children[1].children[1]
        
        #os.makedirs(r'Output')  
        tool_util.write_image('Output/Tool Utilization.png')
        bu_graph.write_image('Output/BU Loading.png')
        turnaround_graph.write_image('Output/TAT.png')
        priority_graph.write_image('Output/High Priority Queue Time Breakdown.png')
        csc_graph.write_image('Output/Overall Analysis Status.png')
        
        type_graph.write_image('Output/Type Breakdown.png')
        type_month_graph.write_image('Output/Type Breakdown by Month.png')
        prod_graph.write_image('Output/Product Breakdown.png')
        
        tech_node.write_image('Output/Technology Node Distribution.png')
        prod_dist.write_image('Output/Technology Product Distribution.png')
        failure_dist.write_image('Output/Technology Failure Distribution.png')


    def create_ppt(self, evt):
        """
        Creates a PowerPoint presentation and adds saved PNG images to the first slide.
        
        Parameters:
        -----------
        evt : event object
            Event object passed to the method.
            
        Returns:
        --------
        None
        """
        self.save_png_tab1()

        X = Presentation('Powerpoint Output/SGP FI_KPI_Template.pptx')
        #Layout = X.slide_layouts[5]
        first_slide = X.slides[0]
        #first_slide.shapes.title.text = ''
        first_slide.shapes.add_picture('Output/TAT.png', Inches(0.2), Inches(1), height = Inches(3))
        first_slide.shapes.add_picture('Output/High Priority Queue Time Breakdown.png', Inches(4.6), Inches(1), \
                                       height = Inches(3))
        first_slide.shapes.add_picture('Output/Overall Analysis Status.png', Inches(9), Inches(1), height = Inches(3))
        first_slide.shapes.add_picture('Output/BU Loading.png', Inches(2.3), Inches(4.2), height = Inches(3))
        first_slide.shapes.add_picture('Output/Tool Utilization.png', Inches(6.9), Inches(4.2), height = Inches(3))

        for i in range(1,6):
            first_slide.shapes[i].line.color.rgb = RGBColor(0, 0, 0)
#         text_frame = first_slide.shapes[0].text_frame

#         p = text_frame.paragraphs[0]
#         p.alignment = PP_ALIGN.LEFT
#         run = p.add_run()
#         run.text = 'SGP FI'

#         font = run.font
#         font.name = 'Calibri'
#         font.size = Pt(30)
#         font.bold = True

        X.save('Powerpoint Output/SGP FI_KPI.pptx')  
        
    def create_ppt_all(self, evt):
        """
        Create a PowerPoint presentation with all the saved images from the dashboard.

        Parameters:
        -----------
        evt : event object
            Event object passed to the method.
            
        Returns:
        --------
        None
        """
        self.save_png_tab1()

        X = Presentation('Powerpoint Output/SGP FI_KPI_Template.pptx')
        #Layout = X.slide_layouts[5]
        first_slide = X.slides[0]
        #first_slide.shapes.title.text = ''
        first_slide.shapes.add_picture('Output/TAT.png', Inches(0.2), Inches(1), height = Inches(3))
        first_slide.shapes.add_picture('Output/High Priority Queue Time Breakdown.png', Inches(4.6), Inches(1), \
                                       height = Inches(3))
        first_slide.shapes.add_picture('Output/Overall Analysis Status.png', Inches(9), Inches(1), height = Inches(3))
        first_slide.shapes.add_picture('Output/BU Loading.png', Inches(2.3), Inches(4.2), height = Inches(3))
        first_slide.shapes.add_picture('Output/Tool Utilization.png', Inches(6.9), Inches(4.2), height = Inches(3))
        
        second_slide = X.slides[1]
        #second_slide.shapes.title.text = ''
        second_slide.shapes.add_picture('Output/Type Breakdown.png', Inches(0.2), Inches(1), height = Inches(3))
        second_slide.shapes.add_picture('Output/Type Breakdown by Month.png', Inches(4.6), Inches(1), height = Inches(3))
        second_slide.shapes.add_picture('Output/Product Breakdown.png', Inches(9), Inches(1), height = Inches(3))
        second_slide.shapes.add_picture('Output/Technology Node Distribution.png', Inches(0.2), Inches(4.2), height = Inches(3))
        second_slide.shapes.add_picture('Output/Technology Product Distribution.png', Inches(4.6), Inches(4.2), \
                                        height = Inches(3))
        second_slide.shapes.add_picture('Output/Technology Product Distribution.png', Inches(9), Inches(4.2), \
                                        height = Inches(3))

        for i in range(1,7):
            if i < 6:
                first_slide.shapes[i].line.color.rgb = RGBColor(0, 0, 0)
            second_slide.shapes[i].line.color.rgb = RGBColor(0, 0, 0)
#         text_frame = first_slide.shapes[0].text_frame

#         p = text_frame.paragraphs[0]
#         p.alignment = PP_ALIGN.LEFT
#         run = p.add_run()
#         run.text = 'SGP FI'

#         font = run.font
#         font.name = 'Calibri'
#         font.size = Pt(30)
#         font.bold = True

        X.save('Powerpoint Output/SGP FI_KPI (ALL).pptx')  
   
     